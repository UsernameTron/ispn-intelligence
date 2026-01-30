#!/usr/bin/env python3
"""
ISPN Board Report Generator v2.0

CRITICAL: This report generator uses ISPN Canonical Calculations ONLY.
All metrics are sourced from ispn_metrics_history.json - NEVER from raw
Genesys exports or legacy kpi_history.json.

Generates weekly and monthly reports in PPTX and DOCX formats using
standardized ISPN formulas from Charlie's LT Scorecard.

Key Differences from v1:
- Shrinkage uses ISPN formula: (Hours Worked - On Queue) / Hours Worked
- Utilization uses FY25 formula with training hours deduction
- Occupancy uses: Call Hours / On Queue Hours
- ACW uses fixed 15 seconds per call

Usage:
    python scripts/board_report.py                          # Generate monthly report
    python scripts/board_report.py --period weekly          # Generate weekly report
    python scripts/board_report.py --format pptx            # PPTX only
    python scripts/board_report.py --month 2025-01          # Specific month
"""

import json
import argparse
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Paths
BASE_DIR = Path(__file__).parent.parent
METRICS_DIR = BASE_DIR / "data" / "metrics"
PARSED_DIR = BASE_DIR / "data" / "parsed"
REPORTS_DIR = BASE_DIR / "reports"
TEMPLATES_DIR = BASE_DIR / "templates"

# Colors for status indicators
COLORS = {
    'GREEN': RGBColor(0x28, 0xA7, 0x45),
    'YELLOW': RGBColor(0xFF, 0xC1, 0x07),
    'RED': RGBColor(0xDC, 0x35, 0x45),
    'DARK': RGBColor(0x1A, 0x1A, 0x2E),
    'ACCENT': RGBColor(0x00, 0xD4, 0xFF),
    'MUTED': RGBColor(0x6C, 0x75, 0x7D),
}

# KPI display configuration
KPI_CONFIG = {
    'fcr': {
        'name': 'First Call Resolution',
        'short': 'FCR',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '> 70%',
        'description': 'Calls resolved without escalation'
    },
    'escalation': {
        'name': 'Escalation Rate',
        'short': 'Escalation',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '< 30%',
        'description': 'Calls requiring escalation'
    },
    'aht': {
        'name': 'Average Handle Time',
        'short': 'AHT',
        'unit': 'min',
        'format': lambda x: f"{x:.2f}" if x else 'N/A',
        'target': '< 10.7 min',
        'description': 'Average call duration including hold and ACW'
    },
    'awt': {
        'name': 'Average Wait Time',
        'short': 'AWT',
        'unit': 'sec',
        'format': lambda x: f"{x:.1f}" if x else 'N/A',
        'target': '< 90 sec',
        'description': 'Average time to answer'
    },
    'shrinkage': {
        'name': 'Shrinkage',
        'short': 'Shrinkage',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '< 30%',
        'description': 'ISPN: (Hours Worked - On Queue) / Hours Worked'
    },
    'utilization': {
        'name': 'Tech Utilization',
        'short': 'Utilization',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '> 55%',
        'description': 'FY25: Inbound Hours / (Hours Worked - Training)'
    },
    'occupancy': {
        'name': 'Occupancy',
        'short': 'Occupancy',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '> 65%',
        'description': 'ISPN: Call Hours / On Queue Hours'
    },
    'quality': {
        'name': 'Quality Score',
        'short': 'Quality',
        'unit': 'pts',
        'format': lambda x: f"{x:.1f}" if x else 'N/A',
        'target': '> 88',
        'description': 'Average tech review score'
    },
    'abandon': {
        'name': 'Abandon Rate',
        'short': 'Abandon',
        'unit': '%',
        'format': lambda x: f"{x:.1%}" if x else 'N/A',
        'target': '< 5%',
        'description': 'Calls abandoned before answer'
    },
}


# =============================================================================
# DATA LOADING - ISPN CANONICAL METRICS ONLY
# =============================================================================

def load_ispn_metrics_history() -> dict:
    """
    Load ISPN-calculated metrics history.
    
    CRITICAL: This is the ONLY valid source for board reports.
    Never use kpi_history.json or raw Genesys values.
    """
    history_file = METRICS_DIR / "ispn_metrics_history.json"
    
    if not history_file.exists():
        print(f"⚠️  ISPN metrics history not found: {history_file}")
        print("   Run: python scripts/ingest.py --calculate")
        return {'periods': {}}
    
    with open(history_file, 'r') as f:
        data = json.load(f)
    
    # Verify it has the expected structure
    if 'periods' not in data:
        print("⚠️  Invalid ISPN metrics history format")
        return {'periods': {}}
    
    return data


def get_latest_period(history: dict) -> Tuple[Optional[str], Optional[dict]]:
    """Get the most recent period and its data."""
    periods = history.get('periods', {})
    if not periods:
        return None, None
    
    latest = max(periods.keys())
    return latest, periods[latest]


def get_previous_period(history: dict, current: str) -> Tuple[Optional[str], Optional[dict]]:
    """Get the period before the current one."""
    periods = sorted(history.get('periods', {}).keys())
    if current in periods:
        idx = periods.index(current)
        if idx > 0:
            prev = periods[idx - 1]
            return prev, history['periods'][prev]
    return None, None


def calculate_deltas(current: dict, previous: dict) -> dict:
    """Calculate changes between periods."""
    if not previous:
        return {}
    
    deltas = {}
    current_kpis = current.get('kpis', {})
    prev_kpis = previous.get('kpis', {})
    
    for kpi, value in current_kpis.items():
        if kpi in prev_kpis and value is not None and prev_kpis[kpi] is not None:
            delta = value - prev_kpis[kpi]
            pct_change = (delta / prev_kpis[kpi] * 100) if prev_kpis[kpi] != 0 else 0
            deltas[kpi] = {
                'current': value,
                'previous': prev_kpis[kpi],
                'delta': delta,
                'pct_change': round(pct_change, 1)
            }
    
    return deltas


# =============================================================================
# NARRATIVE GENERATION
# =============================================================================

def generate_executive_summary(period: str, data: dict, deltas: dict) -> str:
    """Generate executive summary section."""
    statuses = data.get('statuses', {})
    
    # Count by status
    status_counts = {'GREEN': 0, 'YELLOW': 0, 'RED': 0}
    for kpi, info in statuses.items():
        status = info.get('status', 'UNKNOWN')
        if status in status_counts:
            status_counts[status] += 1
    
    total = sum(status_counts.values())
    green_pct = (status_counts['GREEN'] / total * 100) if total > 0 else 0
    
    # Determine overall status
    if status_counts['RED'] >= 2:
        overall = "ATTENTION REQUIRED"
        summary = f"{status_counts['RED']} KPIs in RED status require immediate action."
    elif status_counts['RED'] == 1:
        overall = "MONITORING"
        summary = f"1 KPI in RED status; {status_counts['YELLOW']} in YELLOW."
    elif status_counts['YELLOW'] > 2:
        overall = "MONITORING"
        summary = f"{status_counts['YELLOW']} KPIs in YELLOW status need attention."
    else:
        overall = "ON TRACK"
        summary = f"All {total} KPIs meeting or approaching targets."
    
    return f"""**Overall Status: {overall}**

{summary}

Status Breakdown:
- 🟢 GREEN: {status_counts['GREEN']} ({green_pct:.0f}%)
- 🟡 YELLOW: {status_counts['YELLOW']}
- 🔴 RED: {status_counts['RED']}"""


def generate_kpi_details(data: dict, deltas: dict) -> str:
    """Generate detailed KPI section."""
    statuses = data.get('statuses', {})
    lines = []
    
    for kpi_key, config in KPI_CONFIG.items():
        if kpi_key not in statuses:
            continue
        
        info = statuses[kpi_key]
        value = info.get('value', 'N/A')
        target = info.get('target', 'N/A')
        status = info.get('status', 'UNKNOWN')
        
        # Delta indicator
        delta_info = deltas.get(kpi_key, {})
        delta = delta_info.get('delta')
        
        delta_str = ""
        if delta is not None:
            direction = "↑" if delta > 0 else "↓" if delta < 0 else "→"
            # Determine if direction is good or bad
            if kpi_key in ['fcr', 'utilization', 'occupancy', 'quality']:
                # Higher is better
                color = "improved" if delta > 0 else "declined"
            else:
                # Lower is better
                color = "improved" if delta < 0 else "declined"
            
            if 'pct' in kpi_key or kpi_key == 'fcr':
                delta_str = f" ({direction} {abs(delta):.1%} {color})"
            else:
                delta_str = f" ({direction} {abs(delta):.2f} {color})"
        
        # Status emoji
        emoji = "🟢" if status == "GREEN" else "🟡" if status == "YELLOW" else "🔴"
        
        lines.append(f"- **{config['name']}**: {value} (Target: {target}){delta_str} {emoji}")
        lines.append(f"  *{config['description']}*")
    
    return "\n".join(lines)


def generate_formula_notes() -> str:
    """Generate notes about ISPN formula methodology."""
    return """---

**Calculation Methodology**

All metrics calculated using ISPN canonical formulas from Charlie's LT Scorecard:

| Metric | Formula |
|--------|---------|
| FCR | 1 - (Escalations / Call Tickets) |
| Shrinkage | (Hours Worked - On Queue) / Hours Worked |
| Utilization | Inbound Hours / (Hours Worked - Training) |
| Occupancy | Call Hours / On Queue Hours |
| ACW | Fixed 15 seconds per call |

*Note: These formulas differ from Genesys pre-calculated values.*"""


def generate_narrative(period: str, data: dict, deltas: dict) -> str:
    """Generate complete executive narrative for the report."""
    lines = []
    
    # Header
    lines.append(f"# ISPN Tech Center Performance Report")
    lines.append(f"## {period}")
    lines.append("")
    
    # Metadata
    metadata = data.get('metadata', {})
    if metadata.get('calculation_timestamp'):
        lines.append(f"*Calculated: {metadata['calculation_timestamp'][:19]}*")
        lines.append(f"*Formula Version: {metadata.get('formula_version', '1.0.0')}*")
        lines.append("")
    
    # Executive Summary
    lines.append("## Executive Summary")
    lines.append("")
    lines.append(generate_executive_summary(period, data, deltas))
    lines.append("")
    
    # KPI Details
    lines.append("## KPI Performance")
    lines.append("")
    lines.append(generate_kpi_details(data, deltas))
    lines.append("")
    
    # Call Volume
    call_volume = data.get('call_volume', {})
    if call_volume:
        lines.append("## Call Volume")
        lines.append("")
        lines.append(f"- Inbound Calls: {call_volume.get('inbound_count', 0):,}")
        lines.append(f"- Inbound Hours: {call_volume.get('inbound_hours', 0):,.1f}")
        lines.append(f"- Outbound Calls: {call_volume.get('outbound_count', 0):,}")
        lines.append(f"- Callbacks: {call_volume.get('callback_count', 0):,}")
        lines.append("")
    
    # Agent Hours
    agent_hours = data.get('agent_hours', {})
    if agent_hours:
        lines.append("## Agent Hours")
        lines.append("")
        lines.append(f"- Total Hours Worked: {agent_hours.get('total_hours_worked', 0):,.1f}")
        lines.append(f"- On-Queue Hours: {agent_hours.get('on_queue_hours', 0):,.1f}")
        lines.append(f"- Hours Unavailable (Shrinkage): {agent_hours.get('hours_unavailable', 0):,.1f}")
        lines.append(f"- ACW Hours (ISPN: 15s/call): {agent_hours.get('acw_hours', 0):,.1f}")
        lines.append("")
    
    # Warnings
    warnings = metadata.get('warnings', [])
    if warnings:
        lines.append("## Data Warnings")
        lines.append("")
        for w in warnings:
            lines.append(f"- ⚠️ {w}")
        lines.append("")
    
    # Formula notes
    lines.append(generate_formula_notes())
    lines.append("")
    
    # Footer
    lines.append(f"*Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}*")
    
    return "\n".join(lines)


# =============================================================================
# PPTX GENERATION
# =============================================================================

def create_pptx(period: str, data: dict, deltas: dict, output_path: Path):
    """Generate PowerPoint presentation with ISPN metrics."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "ISPN Tech Center"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = f"Performance Report - {period}"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Formula note
    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.5))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Using ISPN Canonical Calculations (Charlie's LT Scorecard)"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['MUTED']
    p.alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # SLIDE 2: Executive Summary
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Status counts
    statuses = data.get('statuses', {})
    status_counts = {'GREEN': 0, 'YELLOW': 0, 'RED': 0}
    for kpi, info in statuses.items():
        status = info.get('status', 'UNKNOWN')
        if status in status_counts:
            status_counts[status] += 1
    
    total = sum(status_counts.values())
    
    # Status boxes
    y_pos = 1.5
    for status_name, count in [('GREEN', status_counts['GREEN']), 
                                ('YELLOW', status_counts['YELLOW']), 
                                ('RED', status_counts['RED'])]:
        box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(3), Inches(0.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{status_name}: {count} KPIs"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS[status_name]
        p.font.bold = True
        y_pos += 1.0
    
    # Overall status
    if status_counts['RED'] >= 2:
        overall_status = "ATTENTION REQUIRED"
        overall_color = COLORS['RED']
    elif status_counts['RED'] == 1 or status_counts['YELLOW'] > 2:
        overall_status = "MONITORING"
        overall_color = COLORS['YELLOW']
    else:
        overall_status = "ON TRACK"
        overall_color = COLORS['GREEN']
    
    status_box = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(6), Inches(2))
    tf = status_box.text_frame
    p = tf.paragraphs[0]
    p.text = overall_status
    p.font.size = Pt(36)
    p.font.color.rgb = overall_color
    p.font.bold = True
    
    # =========================================================================
    # SLIDE 3: KPI Details
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "KPI Performance (ISPN Canonical)"
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Column headers
    headers = ['KPI', 'Value', 'Target', 'Status', 'Δ']
    col_widths = [3, 2, 2, 1.5, 2]
    x_positions = [0.5]
    for w in col_widths[:-1]:
        x_positions.append(x_positions[-1] + w)
    
    y_pos = 1.2
    for i, (header, x) in enumerate(zip(headers, x_positions)):
        box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(col_widths[i]), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['MUTED']
    
    y_pos = 1.7
    
    # KPI rows
    for kpi_key in ['fcr', 'escalation', 'aht', 'awt', 'shrinkage', 'utilization', 'occupancy', 'quality', 'abandon']:
        if kpi_key not in statuses:
            continue
        
        info = statuses[kpi_key]
        config = KPI_CONFIG.get(kpi_key, {})
        
        value = info.get('value', 'N/A')
        target = info.get('target', 'N/A')
        status = info.get('status', 'UNKNOWN')
        
        delta_info = deltas.get(kpi_key, {})
        delta = delta_info.get('delta')
        
        # KPI name
        box = slide.shapes.add_textbox(Inches(x_positions[0]), Inches(y_pos), Inches(col_widths[0]), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = config.get('short', kpi_key.upper())
        p.font.size = Pt(16)
        
        # Value
        box = slide.shapes.add_textbox(Inches(x_positions[1]), Inches(y_pos), Inches(col_widths[1]), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Target
        box = slide.shapes.add_textbox(Inches(x_positions[2]), Inches(y_pos), Inches(col_widths[2]), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = str(target)
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['MUTED']
        
        # Status
        box = slide.shapes.add_textbox(Inches(x_positions[3]), Inches(y_pos), Inches(col_widths[3]), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "●"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS.get(status, COLORS['MUTED'])
        
        # Delta
        if delta is not None:
            box = slide.shapes.add_textbox(Inches(x_positions[4]), Inches(y_pos), Inches(col_widths[4]), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            direction = "↑" if delta > 0 else "↓" if delta < 0 else "→"
            if 'pct' in kpi_key or kpi_key == 'fcr':
                p.text = f"{direction} {abs(delta):.1%}"
            else:
                p.text = f"{direction} {abs(delta):.2f}"
            p.font.size = Pt(14)
        
        y_pos += 0.55
    
    # =========================================================================
    # SLIDE 4: Formula Methodology
    # =========================================================================
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "ISPN Calculation Methodology"
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Formula table
    formulas = [
        ("FCR", "1 - (Escalations / Call Tickets)", "Helpdesk"),
        ("Shrinkage", "(Hours Worked - On Queue) / Hours Worked", "Agent Status"),
        ("Utilization", "Inbound Hours / (Hours - Training)", "Agent Status"),
        ("Occupancy", "Call Hours / On Queue Hours", "Agent Status"),
        ("ACW", "15 seconds × Call Count (fixed)", "Standard"),
        ("AWT", "Weighted average (Genesys + Wave)", "Interactions"),
    ]
    
    y_pos = 1.5
    for name, formula, source in formulas:
        # Name
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Formula
        box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(7), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = formula
        p.font.size = Pt(14)
        
        # Source
        box = slide.shapes.add_textbox(Inches(10), Inches(y_pos), Inches(3), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[{source}]"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['MUTED']
        
        y_pos += 0.6
    
    # Note
    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Source: Charlie's LT Scorecard Formulas (ISPN_iGLASS_LT_Scorecard_Weekly_Monthly_2025.xlsx)"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = COLORS['MUTED']
    
    # Save
    prs.save(output_path)
    print(f"  ✓ Created: {output_path}")


# =============================================================================
# DOCX GENERATION
# =============================================================================

def create_docx(period: str, narrative: str, output_path: Path):
    """Generate Word document with ISPN metrics."""
    doc = Document()
    
    # Title
    title = doc.add_heading('ISPN Tech Center Performance Report', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_heading(period, level=1)
    
    # Add note about methodology
    p = doc.add_paragraph()
    run = p.add_run("Using ISPN Canonical Calculations (Charlie's LT Scorecard)")
    run.italic = True
    
    doc.add_paragraph()
    
    # Parse narrative and add content
    for line in narrative.split('\n'):
        if line.startswith('# '):
            continue  # Skip main title, we added it above
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('- **'):
            # KPI line with bold portion
            p = doc.add_paragraph(style='List Bullet')
            parts = line[2:].split('**')
            if len(parts) >= 2:
                run = p.add_run(parts[1])
                run.bold = True
                if len(parts) > 2:
                    p.add_run(parts[2])
            else:
                p.add_run(line[2:])
        elif line.startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.startswith('  *') and line.endswith('*'):
            # Description line (indented italic)
            p = doc.add_paragraph()
            run = p.add_run("    " + line[3:-1])
            run.italic = True
            run.font.size = DocxPt(10)
        elif line.startswith('**') and '**' in line[2:]:
            # Bold text
            p = doc.add_paragraph()
            parts = line.split('**')
            for i, part in enumerate(parts):
                if part:
                    run = p.add_run(part)
                    run.bold = (i % 2 == 1)
        elif line.startswith('*') and line.endswith('*'):
            # Italic text
            p = doc.add_paragraph()
            run = p.add_run(line[1:-1])
            run.italic = True
        elif line.startswith('---'):
            doc.add_paragraph('─' * 50)
        elif line.startswith('|'):
            # Table row - skip for now (complex to render)
            continue
        elif line.strip():
            doc.add_paragraph(line)
    
    doc.save(output_path)
    print(f"  ✓ Created: {output_path}")


# =============================================================================
# MAIN ENTRY POINT
# =============================================================================

def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description='ISPN Board Report Generator (ISPN Canonical Calculations)'
    )
    parser.add_argument('--period', choices=['weekly', 'monthly'], default='monthly',
                        help='Report period type')
    parser.add_argument('--month', type=str, help='Month (YYYY-MM) for monthly report')
    parser.add_argument('--week', type=str, help='Week (YYYY-Wxx) for weekly report')
    parser.add_argument('--format', choices=['pptx', 'docx', 'both'], default='both',
                        help='Output format')
    args = parser.parse_args()
    
    print("=" * 60)
    print("ISPN Board Report Generator v2.0")
    print("Using ISPN Canonical Calculations")
    print("=" * 60)
    
    # Load ISPN metrics (NOT legacy kpi_history.json)
    history = load_ispn_metrics_history()
    
    if not history.get('periods'):
        print("\n❌ No ISPN metrics data found.")
        print("   Run: python scripts/ingest.py --calculate")
        print("   to calculate metrics from raw Genesys data.")
        return
    
    # Get period to report on
    if args.month:
        period = args.month
        if period not in history['periods']:
            print(f"\n❌ No data for period: {period}")
            print(f"   Available periods: {', '.join(sorted(history['periods'].keys()))}")
            return
        data = history['periods'][period]
    else:
        period, data = get_latest_period(history)
    
    if not period or not data:
        print("\n❌ No KPI data found.")
        return
    
    print(f"\n📊 Using data from period: {period}")
    
    # Get previous period for deltas
    prev_period, prev_data = get_previous_period(history, period)
    if prev_period:
        print(f"📈 Comparing to previous: {prev_period}")
    else:
        print("   (No previous period for comparison)")
    
    # Calculate deltas
    deltas = calculate_deltas(data, prev_data) if prev_data else {}
    
    # Generate narrative
    narrative = generate_narrative(period, data, deltas)
    
    # Create output directory
    if args.period == 'monthly':
        output_dir = REPORTS_DIR / "board" / period
    else:
        output_dir = REPORTS_DIR / "weekly" / (args.week or period)
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Generate outputs
    print(f"\n📁 Generating reports in: {output_dir}")
    
    # Save narrative markdown
    md_path = output_dir / "narrative.md"
    with open(md_path, 'w') as f:
        f.write(narrative)
    print(f"  ✓ Created: {md_path.name}")
    
    if args.format in ['pptx', 'both']:
        pptx_path = output_dir / "board_report.pptx"
        create_pptx(period, data, deltas, pptx_path)
    
    if args.format in ['docx', 'both']:
        docx_path = output_dir / "board_report.docx"
        create_docx(period, narrative, docx_path)
    
    # Print summary
    print("\n" + "=" * 60)
    print("✅ COMPLETE")
    print("=" * 60)
    print(f"\nReports saved to: {output_dir}")
    
    # Show quick KPI summary
    statuses = data.get('statuses', {})
    print("\nKPI Summary:")
    for kpi_key in ['fcr', 'aht', 'shrinkage', 'utilization']:
        if kpi_key in statuses:
            info = statuses[kpi_key]
            status = info.get('status', 'UNKNOWN')
            emoji = "🟢" if status == "GREEN" else "🟡" if status == "YELLOW" else "🔴"
            print(f"  {emoji} {KPI_CONFIG[kpi_key]['short']}: {info.get('value', 'N/A')}")
    
    # Warnings
    warnings = data.get('metadata', {}).get('warnings', [])
    if warnings:
        print("\n⚠️  Warnings:")
        for w in warnings:
            print(f"   • {w}")
    
    print("\n💡 Next steps:")
    print("  1. Review generated reports")
    print("  2. git add . && git commit -m 'Add board report'")
    print("  3. git push")


if __name__ == '__main__':
    main()
